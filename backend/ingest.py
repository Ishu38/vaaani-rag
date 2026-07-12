"""Document ingestion pipeline.

Two-phase architecture for commercial UX:
  Phase 1 (ingest_vectors, ~1s): embed chunks → write index → save metadata
  Phase 2 (ingest_graph, background): entity extraction → knowledge graph → communities

Reads .pdf/.txt/.md/.docx/.png/.jpg/.jpeg/.webp/.pptx/.html/.htm files.

CLI (full pipeline, synchronous):
    python ingest.py --source ./data/raw --index ./data/index.tq
"""
from __future__ import annotations

import argparse
import hashlib
import json
import os
import sys
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import numpy as np
from sentence_transformers import SentenceTransformer
from turbovec import TurboQuantIndex

from community import build_communities, save_communities
from config import (
    BIT_WIDTH,
    CHUNK_OVERLAP,
    CHUNK_TOKENS,
    COMMUNITIES_PATH,
    DEEPSEEK_TIMEOUT,
    EMBED_DIM,
    EMBED_MODEL_NAME,
    GRAPH_PATH,
    INDEX_PATH,
    METADATA_PATH,
    RAW_DIR,
)
from extractor import extract_chunk
from graph import KnowledgeGraph
import httpx

SUPPORTED_EXT = {
    ".pdf", ".txt", ".md", ".docx",
    ".png", ".jpg", ".jpeg", ".webp",
    ".pptx", ".html", ".htm",
}

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}

MIN_EXTRACT_WORDS = 25
# DeepSeek's API comfortably handles 12-16 concurrent extraction calls
# without rate-limiting on the cheap-chat tier. Bumped from 4 → 12 on
# 2026-05-28 after a commercial-ingest complaint (86-chunk PDF was
# taking ~2.2 minutes on entity extraction alone). At 12 concurrency
# that drops to ~40 s. Env-overridable via VAAANI_EXTRACT_CONCURRENCY
# if a host needs to dial it down (e.g. tight rate-limit tier).
EXTRACT_CONCURRENCY = int(os.environ.get("VAAANI_EXTRACT_CONCURRENCY", "12"))


@dataclass
class Chunk:
    text: str
    source: str
    chunk_no: int


def read_pdf(path: Path) -> str:
    try:
        import fitz
        doc = fitz.open(str(path))
        pages: list[str] = []
        for page in doc:
            pages.append(page.get_text("text") or "")
        doc.close()
        return "\n".join(pages)
    except Exception:
        pass
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        return "\n".join((p.extract_text() or "") for p in reader.pages)
    except Exception:
        raise RuntimeError(f"Failed to read PDF: {path.name}")


def read_docx(path: Path) -> str:
    import docx
    d = docx.Document(str(path))
    return "\n".join(p.text for p in d.paragraphs)


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        if parts:
            slides.append("\n".join(parts))
    return "\n\n".join(slides)


def read_html(path: Path) -> str:
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts: list[str] = []
            self._skip = False

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style", "noscript"):
                self._skip = True

        def handle_endtag(self, tag):
            if tag in ("script", "style", "noscript"):
                self._skip = False
            if tag in ("p", "div", "li", "br", "h1", "h2", "h3", "h4", "h5", "h6", "tr"):
                self.parts.append("\n")

        def handle_data(self, data):
            if not self._skip:
                t = data.strip()
                if t:
                    self.parts.append(t)

    raw = path.read_text(encoding="utf-8", errors="ignore")
    ex = _TextExtractor()
    ex.feed(raw)
    return "\n".join(ex.parts)


def read_image(path: Path) -> str:
    try:
        import pytesseract
    except ImportError:
        raise RuntimeError(
            "pytesseract is required for image ingestion. "
            "Install it: pip install pytesseract && sudo apt install tesseract-ocr"
        )
    from PIL import Image
    img = Image.open(path)
    if path.suffix.lower() == ".webp":
        import io as _io
        buf = _io.BytesIO()
        img = img.convert("RGB")
        img.save(buf, format="JPEG", quality=92)
        buf.seek(0)
        img = Image.open(buf)
    text = pytesseract.image_to_string(img, lang="eng")
    if not text or not text.strip():
        return ""
    return text


def load_document(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return read_pdf(path)
    if ext == ".docx":
        return read_docx(path)
    if ext in {".txt", ".md"}:
        return read_text(path)
    if ext == ".pptx":
        return read_pptx(path)
    if ext in {".html", ".htm"}:
        return read_html(path)
    if ext in IMAGE_EXTS:
        return read_image(path)
    raise ValueError(f"Unsupported file type: {ext}")


def chunk_text(text: str, size: int = CHUNK_TOKENS, overlap: int = CHUNK_OVERLAP) -> list[str]:
    words = text.split()
    if not words:
        return []
    chunks: list[str] = []
    step = max(1, size - overlap)
    for start in range(0, len(words), step):
        window = words[start : start + size]
        if not window:
            break
        chunks.append(" ".join(window))
        if start + size >= len(words):
            break
    return chunks


def file_signature(path: Path) -> str:
    stat = path.stat()
    raw = f"{path.resolve()}|{stat.st_size}|{int(stat.st_mtime)}"
    return hashlib.sha1(raw.encode()).hexdigest()


def load_metadata(path: Path) -> dict:
    if path.exists():
        return json.loads(path.read_text())
    return {"files": {}, "chunks": []}


def save_metadata(meta: dict, path: Path) -> None:
    path.write_text(json.dumps(meta, indent=2, ensure_ascii=False))


def load_or_create_index(index_path: Path) -> tuple[TurboQuantIndex, bool]:
    if index_path.exists():
        try:
            return TurboQuantIndex.load(str(index_path)), True
        except Exception as e:
            print(f"  [warn] failed to load existing index ({e}); starting fresh")
    return TurboQuantIndex(dim=EMBED_DIM, bit_width=BIT_WIDTH), False


def iter_documents(source: Path) -> Iterable[Path]:
    for p in sorted(source.rglob("*")):
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXT:
            yield p


# ═══════════════════════════════════════════════════════════════════════
# Phase 1 — Vectors only (fast, <1-2s for any file)
# ═══════════════════════════════════════════════════════════════════════

def ingest_vectors(source: Path, index_path: Path, metadata_path: Path) -> dict:
    """Embed chunks + write index + save metadata. No entity extraction.

    Returns a dict with the chunks and metadata needed for deferred graph building.
    Caller should reload the retriever after this returns — the file is now
    searchable. Then call ingest_graph_deferred() to build the knowledge graph
    in the background.
    """
    meta = load_metadata(metadata_path)
    index, existed = load_or_create_index(index_path)
    if not existed:
        meta = {"files": {}, "chunks": []}

    model = SentenceTransformer(EMBED_MODEL_NAME)

    total_added = 0
    files_processed = 0
    files_skipped = 0

    t0 = time.time()
    for path in iter_documents(source):
        sig = file_signature(path)
        key = str(path.resolve())
        if meta["files"].get(key, {}).get("signature") == sig:
            files_skipped += 1
            continue

        try:
            text = load_document(path)
        except Exception as e:
            print(f"  [skip] {path.name}: {e}")
            continue

        chunks = chunk_text(text)
        if not chunks:
            print(f"  [skip] {path.name}: empty after chunking")
            continue

        n_chunks = len(chunks)
        est_words = len(text.split())
        print(f"Embedding {path.name} ({n_chunks} chunks, ~{est_words} words)...")

        vectors = model.encode(chunks, show_progress_bar=False, convert_to_numpy=True)
        vectors = np.asarray(vectors, dtype=np.float32)
        index.add(vectors)

        start_idx = len(meta["chunks"])
        for i, ch in enumerate(chunks):
            meta["chunks"].append({
                "source": path.name, "path": key, "chunk_no": i, "text": ch,
            })

        meta["files"][key] = {"signature": sig, "chunks": n_chunks, "name": path.name}
        total_added += n_chunks
        files_processed += 1
        print(f"  {path.name} embedded in {time.time() - t0:.1f}s ({n_chunks} chunks)")

    if total_added > 0:
        index.write(str(index_path))
    save_metadata(meta, metadata_path)

    elapsed = time.time() - t0
    print(f"Vector phase: {elapsed:.1f}s total, {files_processed} new, {files_skipped} skipped")
    return {
        "files_processed": files_processed,
        "files_skipped": files_skipped,
        "chunks_added": total_added,
        "total_chunks": len(meta["chunks"]),
        "triples_added": 0,
        "communities": 0,
        "graph_nodes": 0,
        "graph_edges": 0,
    }


# ═══════════════════════════════════════════════════════════════════════
# Phase 2 — Graph-RAG extraction (background, may take minutes)
# ═══════════════════════════════════════════════════════════════════════

# Chunks per DeepSeek call. 5 is a sweet spot — DeepSeek-chat handles 5
# numbered chunks reliably at temperature 0 with the EXTRACTION_SYSTEM_BATCH
# prompt, blast radius on a single failure is small, and HTTP overhead is
# amortised 5×. Override with VAAANI_EXTRACT_BATCH_SIZE if you want bigger
# batches for cheaper models or smaller for noisier ones.
EXTRACT_BATCH_SIZE = int(os.environ.get("VAAANI_EXTRACT_BATCH_SIZE", "5"))


def _extract_batch_job(
    indices: list[int], texts: list[str], timeout: int
) -> tuple[dict[int, object], str | None]:
    """Run one batched DeepSeek extraction call across `texts`, mapping
    results back onto the corresponding `indices`."""
    from extractor import extract_chunks_batch
    try:
        with httpx.Client(timeout=timeout) as cl:
            batch = extract_chunks_batch(texts, client=cl)
        return {indices[i]: batch[i] for i in range(len(indices))}, None
    except Exception as e:
        return {}, str(e)


def _extract_chunks_parallel(
    chunks: list[str],
    chunk_start_idx: int,
    *,
    concurrency: int = EXTRACT_CONCURRENCY,
    progress_cb=None,
) -> tuple[dict[int, object], int]:
    """Run entity extraction over `chunks` in parallel batches.

    Returns ``(results, failed)`` where `failed` counts chunks whose batch
    errored (timeout, engine down, bad JSON). Callers must NOT mark a file
    graph-extracted when `failed > 0`, otherwise the failure is silent and
    the file never gets a retry — the "zero-node library" bug.

    `progress_cb`, if provided, is called as ``progress_cb(done, total)``
    each time a batch completes — used by main.py to push live progress
    counts into the job row so the SPA can show "extracted X/Y chunks"
    instead of just a phase name.
    """
    results: dict[int, object] = {}
    failed = 0
    jobs: list[tuple[int, str]] = []
    skipped = 0
    for i, ch in enumerate(chunks):
        wc = len(ch.split())
        if wc < MIN_EXTRACT_WORDS:
            skipped += 1
            continue
        jobs.append((chunk_start_idx + i, ch))

    if skipped:
        print(f"  [extract] skipping {skipped}/{len(chunks)} chunks (< {MIN_EXTRACT_WORDS} words)")

    if not jobs:
        return results, failed

    timeout_val = DEEPSEEK_TIMEOUT

    # Group jobs into batches of EXTRACT_BATCH_SIZE.
    batches: list[tuple[list[int], list[str]]] = []
    for b in range(0, len(jobs), EXTRACT_BATCH_SIZE):
        slice_ = jobs[b : b + EXTRACT_BATCH_SIZE]
        batches.append(([idx for idx, _ in slice_], [t for _, t in slice_]))

    total = len(jobs)
    done = 0

    with ThreadPoolExecutor(max_workers=concurrency) as pool:
        futures = {
            pool.submit(_extract_batch_job, idx_list, text_list, timeout_val): (idx_list, text_list)
            for idx_list, text_list in batches
        }
        for future in as_completed(futures):
            batch_result, err = future.result()
            idx_list, _ = futures[future]
            done += len(idx_list)
            if err:
                failed += len(idx_list)
                print(f"  [extract:warn] batch of {len(idx_list)} failed: {err}")
            else:
                for k, v in batch_result.items():
                    if v is not None:
                        results[k] = v
            if progress_cb is not None:
                try:
                    progress_cb(done, total)
                except Exception:
                    pass
            print(f"  [extract] {done}/{total} done ({len(results)} with entities)")

    return results, failed


def ingest_graph_deferred(progress_cb=None) -> dict:
    """Phase 2: entity extraction + knowledge graph + communities.

    Reads the already-saved metadata to find chunks that haven't been
    graph-extracted yet. Runs extraction, saves graph + communities.
    Should be called AFTER ingest_vectors() and retriever.reload().

    `progress_cb`, if supplied, is forwarded to the parallel extractor
    and fires as `progress_cb(done, total)` per completed batch. main.py
    uses this to surface live extraction progress to the SPA.
    """
    meta = load_metadata(METADATA_PATH)
    kg = KnowledgeGraph.load(GRAPH_PATH)

    # Determine which chunks already have graph nodes — we track this by
    # checking if the chunk's entities are already in the graph. Since
    # chunks map 1:1 to extraction results that feed into the graph,
    # we use the metadata "files" table — files whose graph extraction
    # was completed are tracked via graph.json timestamps.
    #
    # Simplest approach: just run extraction for ALL chunks, relying on
    # the kg.ingest_extraction deduplication logic (it upserts nodes/edges).

    # Actually, re-extracting everything is wasteful. Instead, we only
    # extract chunks whose files haven't been graph-processed.
    # We track this with a "_graph_extracted" flag per file in metadata.

    triples_added = 0
    total_chunks = len(meta["chunks"])
    filenames_seen: set[str] = set()

    for file_key, file_info in meta["files"].items():
        if file_info.get("_graph_extracted"):
            continue
        fname = file_info["name"]
        if fname in filenames_seen:
            continue
        filenames_seen.add(fname)

        # Gather all chunks for this file WITH their GLOBAL index into
        # meta["chunks"]. That global index IS the chunk_id the knowledge graph
        # records as provenance and that per-user visibility scoping resolves
        # back to a source file. Using a file-local index here (the old bug)
        # mis-attributed every node after the first file to whatever document
        # occupies that global slot — a cross-document / cross-tenant leak.
        file_chunks: list[tuple[int, str]] = []
        for gi, c in enumerate(meta["chunks"]):
            if c.get("source") == fname:
                file_chunks.append((gi, c["text"]))

        if not file_chunks:
            continue

        n = len(file_chunks)
        local_to_global = [gi for gi, _ in file_chunks]
        print(f"Graph extraction: {fname} ({n} chunks)...")
        t0 = time.time()

        # Normalise every extraction path to key by GLOBAL chunk_id.
        results: dict = {}
        failed_chunks = 0
        if n >= 8:
            # parallel extractor keys by local position (0..n-1); remap to global
            par_results, failed_chunks = _extract_chunks_parallel(
                [t for _, t in file_chunks], 0, progress_cb=progress_cb
            )
            for local_id, ex in par_results.items():
                results[local_to_global[local_id]] = ex
        else:
            # Sequential for small files
            with httpx.Client(timeout=DEEPSEEK_TIMEOUT) as cl:
                for gi, text in file_chunks:
                    wc = len(text.split())
                    if wc < MIN_EXTRACT_WORDS:
                        continue
                    try:
                        results[gi] = extract_chunk(text, client=cl)
                    except Exception as e:
                        failed_chunks += 1
                        print(f"  [extract:warn] {fname} chunk {gi}: {e}")

        ent_count = 0
        rel_count = 0
        for chunk_id, ex in results.items():
            kg.ingest_extraction(ex, chunk_id)
            ent_count += len(ex.entities)
            rel_count += len(ex.relations)
        triples_added += ent_count + rel_count

        # Only mark the file graph-extracted when every extractable chunk
        # succeeded. Marking on failure made extraction errors silent and
        # permanent (a dead engine produced a zero-node graph that never
        # got retried). Leaving the flag unset means the next /ingest run
        # retries just the failed file; already-ingested extractions are
        # deduped by kg.ingest_extraction upserts.
        if failed_chunks == 0:
            meta["files"][file_key]["_graph_extracted"] = True
            save_metadata(meta, METADATA_PATH)
            print(f"  {fname} graph done in {time.time() - t0:.1f}s "
                  f"(+{ent_count} entities, +{rel_count} relations)")
        else:
            print(f"  [extract:warn] {fname}: {failed_chunks}/{n} chunks failed — "
                  f"file left unmarked, will retry on next ingest "
                  f"(+{ent_count} entities, +{rel_count} relations kept)")

    kg.save(GRAPH_PATH)

    communities_count = 0
    if kg.g.number_of_nodes() > 0 and triples_added > 0:
        print("Detecting communities and writing summaries...")
        t0 = time.time()
        communities = build_communities(kg)
        save_communities(communities, COMMUNITIES_PATH)
        communities_count = len(communities)
        print(f"Communities done in {time.time() - t0:.1f}s ({communities_count} communities)")

    result = {
        "files_processed": 0,
        "files_skipped": 0,
        "chunks_added": 0,
        "total_chunks": total_chunks,
        "triples_added": triples_added,
        "communities": communities_count,
        "graph_nodes": kg.g.number_of_nodes(),
        "graph_edges": kg.g.number_of_edges(),
    }
    print(f"Graph phase done. {result}")
    return result


# ═══════════════════════════════════════════════════════════════════════
# Legacy synchronous ingest (CLI: full pipeline)
# ═══════════════════════════════════════════════════════════════════════

def ingest(source: Path, index_path: Path, metadata_path: Path, *, build_graph: bool = True) -> dict:
    """Full synchronous ingest — vectors + graph. For CLI use."""
    result = ingest_vectors(source, index_path, metadata_path)
    if build_graph and result["chunks_added"] > 0:
        graph_result = ingest_graph_deferred()
        result["triples_added"] = graph_result["triples_added"]
        result["communities"] = graph_result["communities"]
        result["graph_nodes"] = graph_result["graph_nodes"]
        result["graph_edges"] = graph_result["graph_edges"]
    return result


def main() -> int:
    ap = argparse.ArgumentParser(description="Ingest documents into the TurboVec RAG index.")
    ap.add_argument("--source", type=Path, default=RAW_DIR, help="Directory with raw documents")
    ap.add_argument("--index", type=Path, default=INDEX_PATH, help="Path to write the TurboVec index")
    ap.add_argument("--metadata", type=Path, default=METADATA_PATH, help="Path to metadata JSON")
    ap.add_argument("--no-graph", action="store_true", help="Skip Graph-RAG extraction (vector-only ingest)")
    args = ap.parse_args()
    ingest(args.source, args.index, args.metadata, build_graph=not args.no_graph)
    return 0


if __name__ == "__main__":
    sys.exit(main())
